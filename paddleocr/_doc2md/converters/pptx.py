# Copyright (c) 2025 PaddlePaddle Authors. All Rights Reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
from pathlib import Path

from ..base import BaseConverter, ConvertResult
from ..registry import default_registry

# pptx XML namespace for DrawingML run properties
_A_STRIKE = "{http://schemas.openxmlformats.org/drawingml/2006/main}strike"
_A_RPR = "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"

# OMML math namespace (same as in DOCX)
_M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_M = "{" + _M_NS + "}"
# DrawingML 2010 extension namespace (PPTX wraps OMML in a14:m)
_A14 = "{http://schemas.microsoft.com/office/drawing/2010/main}"


def _convert_omath(omath_element) -> str:
    """Convert an m:oMath lxml element to LaTeX string. Returns empty string on failure."""
    try:
        from ..math.omml import oMath2Latex

        return str(oMath2Latex(omath_element)).strip()
    except Exception:
        return ""


def _paragraph_has_math(para_element) -> bool:
    """Check if paragraph XML element contains OMML math (a14:m or m:oMath)."""
    return (
        para_element.find(f".//{_A14}m") is not None
        or para_element.find(f".//{_M}oMath") is not None
    )


def _extract_math_from_paragraph(para_element) -> list:
    """Extract LaTeX strings from math elements in a PPTX paragraph XML element."""
    results = []
    # a14:m wraps m:oMathPara or m:oMath
    for a14m in para_element.findall(f".//{_A14}m"):
        for omath in a14m.findall(f".//{_M}oMath"):
            latex = _convert_omath(omath)
            if latex:
                results.append(latex)
        # No oMath inside a14:m? Try the a14:m element itself
        if not results:
            latex = _convert_omath(a14m)
            if latex:
                results.append(latex)
    # Direct m:oMathPara / m:oMath not wrapped in a14:m
    for omath_para in para_element.findall(f".//{_M}oMathPara"):
        for omath in omath_para.findall(f"{_M}oMath"):
            if omath.getparent() is not None and omath.getparent().tag == f"{_A14}m":
                continue  # already handled above
            latex = _convert_omath(omath)
            if latex:
                results.append(latex)
    for omath in para_element.findall(f".//{_M}oMath"):
        parent = omath.getparent()
        if parent is not None and parent.tag in (f"{_A14}m", f"{_M}oMathPara"):
            continue  # already handled
        latex = _convert_omath(omath)
        if latex:
            results.append(latex)
    return results


def _pptx_run_strike(run) -> bool:
    """Return True if the run has strikethrough (sngStrike or dblStrike) set in XML."""
    try:
        rPr = run._r.find(_A_RPR)
        if rPr is not None:
            val = rPr.get(_A_STRIKE)
            return val in ("sngStrike", "dblStrike")
    except Exception:
        pass
    return False


def _pptx_run_script(run) -> str:
    """Return 'super', 'sub', or '' based on DrawingML baseline attribute.

    baseline > 0  -> superscript, baseline < 0 -> subscript.
    baseline attribute is on <a:rPr> with no namespace prefix.
    """
    try:
        rPr = run._r.find(_A_RPR)
        if rPr is not None:
            val = rPr.get("baseline")
            if val is not None:
                n = int(val)
                if n > 0:
                    return "super"
                if n < 0:
                    return "sub"
    except Exception:
        pass
    return ""


def _escape_md_url(url: str) -> str:
    """Escape parentheses in URL for Markdown link syntax."""
    return url.replace("(", "%28").replace(")", "%29")


@default_registry.register
class PptxConverter(BaseConverter):
    supported_extensions = ["pptx"]
    supported_mimetypes = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    def convert_file(self, file_path: Path, **kwargs) -> ConvertResult:
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture  # noqa: F401
        except ImportError:
            raise RuntimeError(
                "PPTX conversion requires python-pptx: pip install paddleocr[doc2md]"
            )

        prs = Presentation(str(file_path))
        slides_md = []
        images: dict = {}
        image_counter = [0]
        slide_width = prs.slide_width

        for slide in prs.slides:
            slide_parts = []

            # Process all shapes
            for shape in slide.shapes:
                self._process_shape(
                    shape, slide_parts, images, image_counter, slide_width, slide.part
                )

            # Handle math formulas inside mc:AlternateContent elements
            # (python-pptx doesn't expose these as Shape objects)
            _MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
            for alt_content in slide._element.iter(f"{{{_MC_NS}}}AlternateContent"):
                # Only look at mc:Choice (the preferred rendering path)
                choice = alt_content.find(f"{{{_MC_NS}}}Choice")
                if choice is None:
                    continue
                for para_elem in choice.iter(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                ):
                    if _paragraph_has_math(para_elem):
                        math_items = _extract_math_from_paragraph(para_elem)
                        for latex in math_items:
                            slide_parts.append(f"$$\n{latex}\n$$")

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Notes**: {notes_text}")

            # Group parts by content type and separate groups with blank lines
            # to prevent HTML blocks from consuming adjacent list items
            def _classify(part: str) -> str:
                s = part.lstrip()
                if s.startswith("##"):
                    return "heading"
                if s.startswith("<"):
                    return "html"
                if s.startswith("- ") or s.lstrip().startswith("- "):
                    return "list"
                if s.startswith(">"):
                    return "blockquote"
                return "other"

            groups: list[list[str]] = []
            for part in slide_parts:
                kind = _classify(part)
                if groups and _classify(groups[-1][0]) == kind:
                    groups[-1].append(part)
                else:
                    groups.append([part])

            slides_md.append("\n\n".join("\n".join(g) for g in groups))

        md_text = "\n\n---\n\n".join(slides_md)

        return ConvertResult(
            markdown=md_text,
            images=images,
            metadata={
                "format": "PPTX",
                "slide_count": len(prs.slides),
            },
        )

    def _process_shape(
        self, shape, slide_parts, images, image_counter, slide_width, slide_part
    ):
        """Recursively process a shape: Picture, GroupShape, Chart, Table, or TextFrame."""
        from pptx.shapes.picture import Picture
        from pptx.util import Emu  # noqa: F401

        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            MSO_SHAPE_TYPE = None

        # 1. Picture
        if isinstance(shape, Picture):
            try:
                img = shape.image
                image_counter[0] += 1
                filename = f"image{image_counter[0]}.{img.ext}"
                rel_path = f"images/{filename}"
                images[rel_path] = img.blob
                if shape.width and slide_width:
                    pct = min(round(shape.width / slide_width * 100), 100)
                    slide_parts.append(f'<img src="images/{filename}" width="{pct}%">')
                else:
                    slide_parts.append(f'<img src="images/{filename}">')
            except (ValueError, AttributeError):
                pass
            return

        # 2. GroupShape - recurse into child shapes
        if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    self._process_shape(
                        child,
                        slide_parts,
                        images,
                        image_counter,
                        slide_width,
                        slide_part,
                    )
            except AttributeError:
                pass
            return

        # 3. Chart
        if shape.has_chart:
            slide_parts.append(self._chart_to_html(shape.chart))
            return

        # 4. Table
        if shape.has_table:
            slide_parts.append(
                self._table_to_html(shape.table, slide_part, image_counter, images)
            )
            return

        # 5. TextFrame
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                # Check for math elements first
                para_xml = paragraph._p
                if _paragraph_has_math(para_xml):
                    math_items = _extract_math_from_paragraph(para_xml)
                    for latex in math_items:
                        slide_parts.append(f"$$\n{latex}\n$$")
                    continue

                parts = []
                for run in paragraph.runs:
                    t = run.text
                    if not t:
                        continue
                    try:
                        url = run.hyperlink.address
                    except Exception:
                        url = None

                    bold = bool(run.font.bold)
                    italic = bool(run.font.italic)
                    underline = bool(run.font.underline) and not url
                    strikethrough = _pptx_run_strike(run)
                    script = _pptx_run_script(run)

                    def _format_segment(
                        seg, bold, italic, underline, strikethrough, script, url
                    ):
                        t = seg
                        if bold or italic or underline or strikethrough or script:
                            leading = len(t) - len(t.lstrip())
                            trailing = len(t) - len(t.rstrip())
                            prefix = t[:leading] if leading else ""
                            suffix = t[len(t) - trailing :] if trailing else ""
                            inner = t.strip()
                            if inner:
                                if strikethrough:
                                    inner = f"~~{inner}~~"
                                if bold and italic:
                                    inner = f"***{inner}***"
                                elif bold:
                                    inner = f"**{inner}**"
                                elif italic:
                                    inner = f"*{inner}*"
                                if underline:
                                    inner = f"<u>{inner}</u>"
                                if script == "super":
                                    inner = f"<sup>{inner}</sup>"
                                elif script == "sub":
                                    inner = f"<sub>{inner}</sub>"
                                t = prefix + inner + suffix
                            elif underline and t:
                                # Pure whitespace + underline = fill-in line
                                # Replace spaces with NBSP so Markdown renderers preserve width
                                t = "<u>" + "\u00a0" * len(t) + "</u>"
                        if url:
                            return f"[{t}]({_escape_md_url(url)})"
                        return t

                    if "\n" in t:
                        segments = t.split("\n")
                        for j, seg in enumerate(segments):
                            if seg:
                                parts.append(
                                    _format_segment(
                                        seg,
                                        bold,
                                        italic,
                                        underline,
                                        strikethrough,
                                        script,
                                        url,
                                    )
                                )
                            if j < len(segments) - 1:
                                parts.append("<br>\n")
                    else:
                        parts.append(
                            _format_segment(
                                t, bold, italic, underline, strikethrough, script, url
                            )
                        )
                text = "".join(parts).strip()
                if not text:
                    continue
                level = paragraph.level
                indent = "  " * level
                # List item continuation lines need indentation; revert <br>\n → <br>
                text = text.replace("<br>\n", "<br>")
                slide_parts.append(f"{indent}- {text}")

    def _chart_to_html(self, chart) -> str:
        """Extract chart data as an HTML table."""
        _CHART_TYPE_NAMES = {
            1: "Area Chart",  # AREA
            2: "Area Chart",  # AREA_STACKED
            4: "Line Chart",  # LINE
            5: "Pie Chart",  # PIE
            15: "Bubble Chart",  # BUBBLE
            51: "Column Chart",  # COLUMN_CLUSTERED
            52: "Column Chart",  # COLUMN_STACKED
            53: "Column Chart",  # COLUMN_STACKED_100
            57: "Bar Chart",  # BAR_CLUSTERED
            58: "Bar Chart",  # BAR_STACKED
            65: "Line Chart",  # LINE_MARKERS
            -4120: "Doughnut Chart",  # DOUGHNUT
            -4151: "Radar Chart",  # RADAR
            -4169: "Scatter Chart",  # XY_SCATTER
        }
        try:
            chart_type_val = chart.chart_type.value if chart.chart_type else 0
            chart_type_name = _CHART_TYPE_NAMES.get(chart_type_val, "Chart")
        except Exception:
            chart_type_name = "Chart"

        try:
            title_text = ""
            try:
                title_text = chart.chart_title.text_frame.text.strip()
            except Exception:
                pass

            # Extract axis info from OOXML
            C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            _C = "{" + C_NS + "}"
            _A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            chart_root = chart._element

            cat_ax_title = ""
            for ax_tag in (f"{_C}catAx", f"{_C}dateAx"):
                ax_el = chart_root.find(f".//{ax_tag}")
                if ax_el is not None:
                    t_el = ax_el.find(f"{_C}title")
                    if t_el is not None:
                        texts = t_el.findall(f".//{_A}t")
                        cat_ax_title = "".join(el.text or "" for el in texts).strip()
                    break

            has_date_ax = chart_root.find(f".//{_C}dateAx") is not None

            val_ax_title = ""
            val_ax_el = chart_root.find(f".//{_C}valAx")
            if val_ax_el is not None:
                t_el = val_ax_el.find(f"{_C}title")
                if t_el is not None:
                    texts = t_el.findall(f".//{_A}t")
                    val_ax_title = "".join(el.text or "" for el in texts).strip()

            plot = chart.plots[0]
            categories = list(plot.categories) if plot.categories else []
            series_list = list(plot.series)

            if not series_list:
                return f"[{chart_type_name}]"

            # Convert Excel date serials to YYYY-MM-DD for date axes
            if has_date_ax and categories:
                from datetime import datetime, timedelta

                converted = []
                for c in categories:
                    try:
                        dt = datetime(1899, 12, 30) + timedelta(days=float(c))
                        converted.append(dt.strftime("%Y-%m-%d"))
                    except (ValueError, TypeError):
                        converted.append(str(c) if c is not None else "")
                categories = converted

            # Collect series names and values
            series_names = []
            series_values = []
            for idx, series in enumerate(series_list):
                try:
                    name = (
                        (series.tx.text if series.tx else "")
                        or val_ax_title
                        or f"Series{idx+1}"
                    )
                except Exception:
                    name = val_ax_title or f"Series{idx+1}"
                series_names.append(name)
                try:
                    vals = [
                        str(round(v, 4)) if v is not None else "" for v in series.values
                    ]
                except Exception:
                    vals = []
                series_values.append(vals)

            # Build HTML table
            html_parts = ["<table>"]
            if title_text:
                html_parts.append(f"<caption>{title_text}</caption>")

            has_header = cat_ax_title or any(name for name in series_names)
            if has_header:
                html_parts.append("<thead><tr>")
                html_parts.append(f"<th>{cat_ax_title}</th>")
                for name in series_names:
                    html_parts.append(f"<th>{name}</th>")
                html_parts.append("</tr></thead>")

            html_parts.append("<tbody>")
            if categories:
                for i, cat in enumerate(categories):
                    html_parts.append(f"<tr><td>{cat}</td>")
                    for vals in series_values:
                        v = vals[i] if i < len(vals) else ""
                        html_parts.append(f"<td>{v}</td>")
                    html_parts.append("</tr>")
            else:
                max_len = max((len(v) for v in series_values), default=0)
                for i in range(max_len):
                    html_parts.append(f"<tr><td>Item{i+1}</td>")
                    for vals in series_values:
                        v = vals[i] if i < len(vals) else ""
                        html_parts.append(f"<td>{v}</td>")
                    html_parts.append("</tr>")
            html_parts.append("</tbody></table>")

            return "\n".join(html_parts)
        except Exception:
            return f"[{chart_type_name}]"

    @staticmethod
    def _table_to_html(
        table, slide_part, image_counter_list: list, images: dict
    ) -> str:
        """Convert a PPTX table to an HTML table, handling merged cells and cell background images."""
        _BLIP_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        _REL_NS = (
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        )

        visited: set[tuple[int, int]] = set()
        html_parts = ["<table>"]

        for i, row in enumerate(table.rows):
            html_parts.append("<tr>")
            for j, cell in enumerate(row.cells):
                if (i, j) in visited:
                    continue
                tag = "th" if i == 0 else "td"
                attrs = ""
                if cell.is_merge_origin:
                    rs = cell.span_height
                    cs = cell.span_width
                    if cs > 1:
                        attrs += f' colspan="{cs}"'
                    if rs > 1:
                        attrs += f' rowspan="{rs}"'
                    for di in range(rs):
                        for dj in range(cs):
                            if (di, dj) != (0, 0):
                                visited.add((i + di, j + dj))

                content_parts = []

                # Extract cell background blip images
                blips = cell._tc.findall(f".//{_BLIP_NS}blip")
                for blip in blips:
                    r_embed = blip.get(f"{_REL_NS}embed")
                    if r_embed:
                        try:
                            image_part = slide_part.related_parts[r_embed]
                            image_counter_list[0] += 1
                            ext = image_part.content_type.split("/")[-1]
                            filename = f"image{image_counter_list[0]}.{ext}"
                            rel_path = f"images/{filename}"
                            images[rel_path] = image_part.blob
                            content_parts.append(
                                f'<img src="images/{filename}" width="100%">'
                            )
                        except (KeyError, AttributeError):
                            pass

                cell_text_parts = []
                for para in cell.text_frame.paragraphs:
                    # Check for math elements first
                    para_xml = para._p
                    if _paragraph_has_math(para_xml):
                        math_items = _extract_math_from_paragraph(para_xml)
                        for latex in math_items:
                            cell_text_parts.append(f"${latex}$")
                        continue

                    run_parts = []
                    for run in para.runs:
                        t = run.text or ""
                        if not t:
                            continue
                        try:
                            url = run.hyperlink.address
                        except Exception:
                            url = None

                        bold = bool(run.font.bold)
                        italic = bool(run.font.italic)
                        underline = bool(run.font.underline) and not url
                        strikethrough = bool(run.font.strike)
                        script = _pptx_run_script(run)

                        if bold:
                            t = f"<b>{t}</b>"
                        if italic:
                            t = f"<i>{t}</i>"
                        if underline:
                            t = f"<u>{t}</u>"
                        if strikethrough:
                            t = f"<del>{t}</del>"
                        if script == "super":
                            t = f"<sup>{t}</sup>"
                        elif script == "sub":
                            t = f"<sub>{t}</sub>"

                        if url:
                            run_parts.append(f'<a href="{url}">{t}</a>')
                        else:
                            run_parts.append(t)
                    cell_text_parts.append("".join(run_parts))
                text = "<br>".join(p for p in cell_text_parts if p.strip())
                if text:
                    content_parts.append(text)
                cell_html = "<br>".join(content_parts) if content_parts else ""

                html_parts.append(f"<{tag}{attrs}>{cell_html}</{tag}>")
            html_parts.append("</tr>")

        html_parts.append("</table>")
        return "\n".join(html_parts)
